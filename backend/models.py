def get_teaching_style(grade):
    if grade <= 5:
        return "Activity-based learning, storytelling, visual aids"
    elif grade <= 8:
        return "Concept explanation + group activities"
    elif grade <= 12:
        return "Discussion-based, analytical thinking, problem solving"
    else:
        return "Advanced critical thinking and project-based learning"


import os
import re
from typing import Optional, List, Dict


def _load_source_text(path: str) -> Optional[str]:
    """Attempt to load plain text from a file path. Returns None if unreadable.

    Supports common text-based files (.txt, .md, .py, .json, .csv, .html).
    """
    if not path:
        return None
    try:
        ext = os.path.splitext(path)[1].lower()
        if ext in (".txt", ".md", ".py", ".json", ".csv", ".html"):
            with open(path, "r", encoding="utf-8") as f:
                return f.read()

        if ext == ".pdf":
            try:
                from PyPDF2 import PdfReader

                reader = PdfReader(path)
                pages = []
                for p in reader.pages:
                    text = p.extract_text()
                    if text:
                        pages.append(text)
                return "\n\n".join(pages)
            except Exception:
                return None

        if ext == ".docx":
            try:
                from docx import Document

                doc = Document(path)
                paras = [p.text for p in doc.paragraphs if p.text]
                return "\n\n".join(paras)
            except Exception:
                return None

        if ext == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt = shape.text.strip()
                            if txt:
                                texts.append(txt)
                return "\n\n".join(texts)
            except Exception:
                return None

        # Fallback: try to open as text
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return None


def _summarize_source(text: Optional[str], max_chars: int = 400) -> str:
    if not text:
        return ""
    text = text.strip()
    # take first paragraph or first max_chars
    parts = re.split(r"\n\s*\n", text)
    if parts and parts[0]:
        return parts[0][:max_chars]
    return text[:max_chars]


def _extract_key_terms(topic: str, source: str, max_terms: int = 6) -> List[str]:
    text = (topic or "") + " " + (source or "")
    # simple tokenization + stopwords removal
    stopwords = set(
        "a an the and or but if then when where how to of for on in with by at from is are was were be been".split()
    )
    tokens = re.findall(r"[A-Za-z]{3,}", text)
    freq: Dict[str, int] = {}
    for t in tokens:
        w = t.lower()
        if w in stopwords:
            continue
        freq[w] = freq.get(w, 0) + 1
    terms = sorted(freq.items(), key=lambda x: x[1], reverse=True)
    return [t[0] for t in terms[:max_terms]]


def _make_objectives(topic: str, grade: int, key_terms: List[str]) -> List[str]:
    verbs = ["Identify", "Explain", "Apply", "Analyze", "Create"]
    level = 0 if grade <= 5 else (1 if grade <= 8 else (2 if grade <= 12 else 3))
    obj_count = min(4, 3 + level)
    objs = []
    for i in range(obj_count):
        verb = verbs[min(i + level, len(verbs) - 1)]
        term = key_terms[i] if i < len(key_terms) else topic
        objs.append(f"{verb} the concept of {term} in everyday contexts.")
    return objs


def _make_quiz(topic: str, key_terms: List[str], grade: int, num_questions: int = 4):
    questions = []
    base = key_terms[:]
    if not base:
        base = topic.split()[:3]
    for i in range(num_questions):
        term = base[i % len(base)]
        q = {
            "question": f"Which statement best describes '{term}'?",
            "options": [
                f"A brief definition of {term}.",
                f"A related but incorrect interpretation of {term}.",
                f"An example of {term} in practice.",
                f"An unrelated concept to distract learners.",
            ],
            "answer_index": 0,
            "explanation": f"'{term}' refers to the idea described in option A; options B-D are distractors for formative assessment.",
        }
        questions.append(q)
    return questions


def generate_lesson_plan(
    topic: str,
    grade: int = 6,
    source_text: Optional[str] = None,
    source_file_path: Optional[str] = None,
    duration_minutes: int = 40,
) -> dict:
    """Generate a structured lesson plan (default 40 minutes) from a topic or source.

    This function preserves existing logic in `get_teaching_style` and builds
    a deterministic, pedagogically-structured plan suitable for classrooms.
    It does not call external APIs; it uses simple heuristics to create goals,
    activities, interactive exercises, and a short quiz.
    """
    # Load source if path provided
    if source_file_path and not source_text:
        source_text = _load_source_text(source_file_path)

    source_summary = _summarize_source(source_text)
    key_terms = _extract_key_terms(topic, source_summary)

    teaching_style = get_teaching_style(grade)

    objectives = _make_objectives(topic, grade, key_terms)

    # Simple time breakdown for the lesson (sum exactly to duration_minutes)
    intro = max(4, round(duration_minutes * 0.12))
    warmup = max(4, round(duration_minutes * 0.12))
    instruction = max(15, round(duration_minutes * 0.4))
    interactive = max(6, round(duration_minutes * 0.2))
    quiz_time = max(3, round(duration_minutes * 0.08))
    used = intro + warmup + instruction + interactive + quiz_time
    wrap = max(1, duration_minutes - used)

    segments = [
        {"name": "Introduction", "minutes": intro, "purpose": f"Hook learners and present the learning objectives about {topic}."},
        {"name": "Warm-up", "minutes": warmup, "purpose": "Activate prior knowledge with a short prompt or question."},
        {"name": "Direct Instruction", "minutes": instruction, "purpose": "Teacher-led explanation and modelling of the core concept(s)."},
        {"name": "Interactive Exercise", "minutes": interactive, "purpose": "Hands-on or group activity to practice the concept."},
        {"name": "Quiz", "minutes": quiz_time, "purpose": "Short formative assessment to check understanding."},
        {"name": "Wrap-up", "minutes": wrap, "purpose": "Summarize and set follow-up tasks or homework."},
    ]

    # Generate activity descriptions
    activities = []
    activities.append({
        "title": "Hook: Quick Scenario",
        "duration": segments[0]["minutes"],
        "description": f"Present a relatable scenario that illustrates {topic}. Ask students: 'What would you do?'.",
        "materials": ["Whiteboard", "Marker"],
    })

    activities.append({
        "title": "Think-Pair-Share",
        "duration": segments[1]["minutes"],
        "description": f"Students discuss prior experiences related to {topic} in pairs and share with the class.",
        "materials": ["Prompt handout"],
    })

    activities.append({
        "title": "Guided Practice",
        "duration": segments[2]["minutes"],
        "description": f"Teacher explains core ideas and works through 2-3 examples drawing from {', '.join(key_terms[:3]) or topic}.",
        "materials": ["Slide deck", "Examples worksheet"],
    })

    activities.append({
        "title": "Interactive Group Task",
        "duration": segments[3]["minutes"],
        "description": f"Students work in small groups to complete a hands-on task applying {topic}. Include roles and a success criteria.",
        "materials": ["Task sheet", "Station materials"],
    })

    quiz = _make_quiz(topic, key_terms, grade, num_questions= max(3, min(5, len(key_terms) or 3)))

    plan = {
        "title": f"{duration_minutes}-minute Lesson: {topic}",
        "topic": topic,
        "grade": grade,
        "duration_minutes": duration_minutes,
        "teaching_style": teaching_style,
        "objectives": objectives,
        "key_terms": key_terms,
        "source_summary": source_summary,
        "segments": segments,
        "activities": activities,
        "quiz": quiz,
        "differentiation": {
            "support": "Provide sentence starters, visual aids, and worked examples.",
            "extension": "Challenge students with an open-ended task or project to deepen understanding.",
        },
        "formative_assessment": "Use exit ticket questions and quick checks for understanding during activities.",
    }

    return plan


__all__ = ["get_teaching_style", "generate_lesson_plan"]